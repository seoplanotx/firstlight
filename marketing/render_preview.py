"""Lightweight pptx -> PNG previewer (for visual QA only).

Reads shapes/fills/text from the generated deck and draws an approximation
with Pillow. Not a full renderer; good enough to catch overlaps/legibility.
"""
import os
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

HERE = os.path.dirname(__file__)
SCALE = 100  # px per inch
EMU_PER_IN = 914400


def px(emu):
    return int(round(emu / EMU_PER_IN * SCALE))


FONTS = {
    "Arial": {
        False: "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        True: "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    },
    "Georgia": {
        False: "/usr/share/fonts/truetype/dejavu/DejaVuSerif.ttf",
        True: "/usr/share/fonts/truetype/dejavu/DejaVuSerif-Bold.ttf",
    },
}
_cache = {}


def font(name, size_pt, bold):
    path = FONTS.get(name, FONTS["Arial"]).get(bold)
    key = (path, size_pt)
    if key not in _cache:
        _cache[key] = ImageFont.truetype(path, int(size_pt * SCALE / 72))
    return _cache[key]


def rgb(c):
    return (c[0], c[1], c[2])


def slide_bg(slide):
    try:
        f = slide.background.fill
        if f.type is not None and f.fore_color and f.fore_color.rgb:
            return rgb(f.fore_color.rgb)
    except Exception:
        pass
    return (255, 255, 255)


def shape_fill(sp):
    try:
        if sp.fill.type is not None:
            return rgb(sp.fill.fore_color.rgb)
    except Exception:
        return None
    return None


def wrap(draw, words_runs, fnt, maxw):
    """wrap a single paragraph's text (joined) to width; returns list of lines."""
    text = "".join(t for t, _ in words_runs)
    words = text.split(" ")
    lines, cur = [], ""
    for w in words:
        trial = (cur + " " + w).strip()
        if draw.textlength(trial, font=fnt) <= maxw or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


def render_slide(slide, idx):
    W = px(prs.slide_width)
    H = px(prs.slide_height)
    img = Image.new("RGB", (W, H), slide_bg(slide))
    d = ImageDraw.Draw(img)
    for sp in slide.shapes:
        x, y, w, h = px(sp.left or 0), px(sp.top or 0), px(sp.width or 0), px(sp.height or 0)
        if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                import io as _io
                pic = Image.open(_io.BytesIO(sp.image.blob)).convert("RGB").resize((max(1, w), max(1, h)))
                img.paste(pic, (x, y))
            except Exception:
                pass
            continue
        fill = shape_fill(sp)
        is_auto = sp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        # draw shape background
        if fill is not None and (is_auto):
            try:
                name = sp.adjustments  # noqa
            except Exception:
                pass
            shp = None
            try:
                shp = sp.auto_shape_type
            except Exception:
                shp = None
            sname = str(shp)
            if "OVAL" in sname:
                d.ellipse([x, y, x + w, y + h], fill=fill)
            elif "ROUNDED" in sname:
                r = min(18, w // 2, h // 2)
                d.rounded_rectangle([x, y, x + w, y + h], radius=r, fill=fill)
            elif "HEART" in sname:
                d.ellipse([x, y, x + w, y + h], fill=fill)
            else:
                d.rectangle([x, y, x + w, y + h], fill=fill)
        # draw text
        if sp.has_text_frame and sp.text_frame.text.strip():
            tf = sp.text_frame
            ml = px(tf.margin_left or 0) or 6
            mr = px(tf.margin_right or 0) or 6
            tx0 = x + ml
            tmaxw = w - ml - mr
            # measure total height
            paras = []
            for para in tf.paragraphs:
                runs = [(r.text, r) for r in para.runs if r.text]
                if not runs:
                    # empty spacer
                    sz = para.runs[0].font.size.pt if (para.runs and para.runs[0].font.size) else 10
                    paras.append(("spacer", sz, None, None, 0))
                    continue
                r0 = runs[0][1]
                sz = (r0.font.size.pt if r0.font.size else 14)
                bold = bool(r0.font.bold)
                fname = r0.font.name or "Arial"
                col = rgb(r0.font.color.rgb) if (r0.font.color and r0.font.color.type is not None) else (20, 20, 20)
                fnt = font(fname, sz, bold)
                lines = wrap(d, [(t, rr) for t, rr in runs], fnt, tmaxw)
                ls = (para.line_spacing or 1.0)
                lineh = sz * SCALE / 72 * 1.2 * ls
                align = para.alignment
                paras.append(("text", sz, fnt, (lines, col, lineh, align), 0))
            total_h = 0
            for kind, sz, fnt, data, _ in paras:
                if kind == "spacer":
                    total_h += sz * SCALE / 72
                else:
                    lines, col, lineh, align = data
                    total_h += lineh * len(lines)
            # vertical anchor
            anchor = tf.vertical_anchor
            if anchor == MSO_ANCHOR.MIDDLE:
                cy = y + (h - total_h) / 2
            elif anchor == MSO_ANCHOR.BOTTOM:
                cy = y + h - total_h
            else:
                cy = y + (px(tf.margin_top or 0) or 4)
            for kind, sz, fnt, data, _ in paras:
                if kind == "spacer":
                    cy += sz * SCALE / 72
                    continue
                lines, col, lineh, align = data
                for ln in lines:
                    lw = d.textlength(ln, font=fnt)
                    if align == PP_ALIGN.CENTER:
                        lx = tx0 + (tmaxw - lw) / 2
                    elif align == PP_ALIGN.RIGHT:
                        lx = tx0 + (tmaxw - lw)
                    else:
                        lx = tx0
                    d.text((lx, cy), ln, font=fnt, fill=col)
                    cy += lineh
    return img


prs = Presentation(os.path.join(HERE, "Coffey-Promotional-Deck.pptx"))
imgs = []
for i, slide in enumerate(prs.slides, 1):
    im = render_slide(slide, i)
    imgs.append(im)

# contact sheet: 3 cols
cols = 3
rows = (len(imgs) + cols - 1) // cols
tw, th = imgs[0].size
pad = 24
sheet = Image.new("RGB", (cols * tw + (cols + 1) * pad, rows * th + (rows + 1) * pad), (230, 233, 236))
for i, im in enumerate(imgs):
    r, c = divmod(i, cols)
    sheet.paste(im, (pad + c * (tw + pad), pad + r * (th + pad)))
sheet = sheet.resize((sheet.width // 2, sheet.height // 2))
sheet.save(os.path.join(HERE, "preview-contact-sheet.png"))
print("wrote contact sheet", sheet.size)
