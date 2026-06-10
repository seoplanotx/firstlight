"""
Firstlight — promotional deck generator.

Produces a 16:9 .pptx that imports cleanly into Google Slides
(upload to Drive -> "Open with Google Slides").

Run:
    pip install python-pptx
    python marketing/build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
INK      = RGBColor(0x10, 0x2A, 0x43)   # deep navy ink
NAVY     = RGBColor(0x0B, 0x20, 0x35)   # darker hero navy
TEAL     = RGBColor(0x0E, 0x7C, 0x86)   # primary brand teal
TEAL_BR  = RGBColor(0x14, 0xB8, 0xA6)   # bright teal accent
CORAL    = RGBColor(0xE8, 0x77, 0x5A)   # warm human accent
AMBER    = RGBColor(0xF4, 0xA2, 0x59)   # warm secondary
LIGHT    = RGBColor(0xEF, 0xF5, 0xF6)   # light panel
CLOUD    = RGBColor(0xF7, 0xFA, 0xFB)   # near-white panel
MUTED    = RGBColor(0x5B, 0x6B, 0x7A)   # muted text
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SLATE    = RGBColor(0x33, 0x44, 0x55)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

FONT = "Arial"
SERIF = "Georgia"


# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font, *rest) in para:
            italic = rest[0] if rest else False
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
            r.font.italic = italic
    return tb


def chip(s, x, y, txt, fill, fg=WHITE, w=None, size=11):
    w = w or Inches(0.18 + 0.092 * len(txt))
    sp = rect(s, x, y, w, Inches(0.34), fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        sp.adjustments[0] = 0.5
    except Exception:
        pass
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    r.font.name = FONT
    return sp


def kicker(s, x, y, txt, color=TEAL_BR):
    return text(s, x, y, Inches(8), Inches(0.4),
                [[(txt.upper(), 13, color, True, FONT)]])


def page_tag(s, idx, label):
    text(s, Inches(11.7), Inches(7.02), Inches(1.5), Inches(0.34),
         [[(f"Firstlight  ·  {idx:02d}", 9, MUTED, False, FONT)]],
         align=PP_ALIGN.RIGHT)
    rect(s, Inches(0.0), Inches(7.34), EMU_W, Inches(0.16), TEAL)


def card(s, x, y, w, h, fill=WHITE):
    c = rect(s, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        c.adjustments[0] = 0.06
    except Exception:
        pass
    return c


def logo_mark(s, x, y, scale=1.0):
    """Simple shield/pulse brand mark."""
    sz = Inches(0.62 * scale)
    sh = s.shapes.add_shape(MSO_SHAPE.HEART, x, y, sz, sz)
    sh.fill.solid()
    sh.fill.fore_color.rgb = CORAL
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# ================================================================ SLIDE 1 — HERO
s = slide()
bg(s, NAVY)
# layered accent panels
rect(s, Inches(0), Inches(0), Inches(0.32), EMU_H, TEAL)
rect(s, Inches(0.32), Inches(0), Inches(0.10), EMU_H, TEAL_BR)
# soft right panel
rect(s, Inches(9.2), Inches(0), Inches(4.13), EMU_H, RGBColor(0x0E, 0x2A, 0x42))

logo_mark(s, Inches(1.0), Inches(1.0), scale=1.15)
text(s, Inches(1.75), Inches(1.02), Inches(7), Inches(0.7),
     [[("Firstlight", 26, WHITE, True, FONT)]])

text(s, Inches(1.0), Inches(2.55), Inches(9.6), Inches(2.2),
     [[("Never miss what", 46, WHITE, True, FONT)],
      [("might matter.", 46, TEAL_BR, True, FONT)]],
     line_spacing=1.02)

text(s, Inches(1.02), Inches(4.55), Inches(8.4), Inches(1.3),
     [[("A local-first companion that watches oncology research for you — ",
        18, RGBColor(0xC9, 0xD8, 0xE0), False, FONT)],
      [("and turns it into source-backed notes for your care team.",
        18, RGBColor(0xC9, 0xD8, 0xE0), False, FONT)]],
     line_spacing=1.2)

chip(s, Inches(1.02), Inches(5.95), "Private by design", TEAL, WHITE)
chip(s, Inches(3.05), Inches(5.95), "Source-backed", TEAL_BR, NAVY)
chip(s, Inches(4.85), Inches(5.95), "Built for families", CORAL, WHITE)

# right-side vertical accent text
text(s, Inches(9.55), Inches(2.7), Inches(3.4), Inches(2.4),
     [[("CLINICALTRIALS.GOV", 12, TEAL_BR, True, FONT)],
      [("PUBMED", 12, TEAL_BR, True, FONT)],
      [("", 6, WHITE, False, FONT)],
      [("Monitored daily.", 15, WHITE, True, FONT)],
      [("Matched to one profile.", 15, RGBColor(0xC9,0xD8,0xE0), False, FONT)]],
     line_spacing=1.25)
rect(s, Inches(0.0), Inches(7.34), EMU_W, Inches(0.16), TEAL_BR)


# ================================================================ SLIDE 2 — THE PROBLEM
s = slide()
bg(s, CLOUD)
rect(s, Inches(0), Inches(0), EMU_W, Inches(1.55), WHITE)
kicker(s, Inches(0.9), Inches(0.55), "The problem", CORAL)
text(s, Inches(0.9), Inches(0.86), Inches(11.5), Inches(0.9),
     [[("Hope moves fast. Families can't keep up alone.", 30, INK, True, FONT)]])

stats = [
    ("~1,500+", "new oncology studies indexed every week", TEAL),
    ("Thousands", "of cancer trials recruiting at any moment", TEAL_BR),
    ("Minutes", "the average appointment leaves for new research", CORAL),
]
x = Inches(0.9)
for big, small, col in stats:
    card(s, x, Inches(2.0), Inches(3.75), Inches(2.0), WHITE)
    rect(s, x, Inches(2.0), Inches(3.75), Inches(0.14), col,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.3), Inches(2.35), Inches(3.2), Inches(0.8),
         [[(big, 30, col, True, FONT)]])
    text(s, x + Inches(0.3), Inches(3.15), Inches(3.2), Inches(0.8),
         [[(small, 14, SLATE, False, FONT)]], line_spacing=1.12)
    x = x + Inches(4.0)

text(s, Inches(0.9), Inches(4.45), Inches(11.5), Inches(2.2),
     [[("When someone you love is diagnosed, the breakthroughs are out there — ",
        18, INK, True, FONT)],
      [("buried in journals and registries, written for clinicians, scattered across "
        "sources, and changing by the week.",
        18, SLATE, False, FONT)],
      [("", 8, SLATE, False, FONT)],
      [("The information that could matter most often surfaces too late, or never reaches the right conversation.",
        18, CORAL, True, FONT)]],
     line_spacing=1.22)
page_tag(s, 2, "")


# ================================================================ SLIDE 3 — MEET ONCOWATCH
s = slide()
bg(s, NAVY)
rect(s, Inches(0), Inches(0), Inches(0.32), EMU_H, TEAL_BR)
kicker(s, Inches(0.9), Inches(0.7), "Meet Firstlight", TEAL_BR)
text(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.4),
     [[("Your tireless research assistant — ", 34, WHITE, True, FONT)],
      [("working quietly on your kitchen counter.", 34, TEAL_BR, True, FONT)]],
     line_spacing=1.05)

feats = [
    ("Watches", "Scans ClinicalTrials.gov and PubMed using your loved one's profile."),
    ("Understands", "Matches and scores findings with transparent, deterministic rules."),
    ("Explains", "Plain-language summaries with the reason each item surfaced."),
    ("Prepares", "Exports a clean report to bring to your oncology team."),
]
x = Inches(0.9)
for title, body in feats:
    card(s, x, Inches(3.05), Inches(2.78), Inches(3.0), RGBColor(0x12,0x32,0x4D))
    rect(s, x + Inches(0.32), Inches(3.4), Inches(0.5), Inches(0.5), TEAL_BR,
         shape=MSO_SHAPE.OVAL)
    text(s, x + Inches(0.32), Inches(4.1), Inches(2.2), Inches(0.5),
         [[(title, 19, WHITE, True, FONT)]])
    text(s, x + Inches(0.32), Inches(4.62), Inches(2.25), Inches(1.4),
         [[(body, 13.5, RGBColor(0xC2,0xD2,0xDC), False, FONT)]], line_spacing=1.18)
    x = x + Inches(2.97)
rect(s, Inches(0.0), Inches(7.34), EMU_W, Inches(0.16), TEAL_BR)


# ================================================================ SLIDE 4 — HOW IT WORKS
s = slide()
bg(s, CLOUD)
kicker(s, Inches(0.9), Inches(0.55), "How it works", TEAL)
text(s, Inches(0.9), Inches(0.86), Inches(11.5), Inches(0.9),
     [[("Four quiet steps. No terminal. No cloud required.", 30, INK, True, FONT)]])

steps = [
    ("1", "Build the profile", "Cancer type, subtype, biomarkers, stage, and prior therapies — entered once.", TEAL),
    ("2", "Monitor sources", "Firstlight queries trusted public sources on a schedule while it's open.", TEAL_BR),
    ("3", "Match & score", "Deterministic rules rank what's relevant and flag what's missing.", AMBER),
    ("4", "Brief your team", "Export a source-backed PDF with questions for your oncologist.", CORAL),
]
y = Inches(2.05)
for num, title, body, col in steps:
    card(s, Inches(0.9), y, Inches(11.5), Inches(1.06), WHITE)
    circ = rect(s, Inches(1.15), y + Inches(0.22), Inches(0.62), Inches(0.62),
                col, shape=MSO_SHAPE.OVAL)
    tf = circ.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = FONT
    text(s, Inches(2.05), y + Inches(0.16), Inches(3.3), Inches(0.8),
         [[(title, 19, INK, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5.3), y + Inches(0.16), Inches(6.8), Inches(0.8),
         [[(body, 15, SLATE, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
    y = y + Inches(1.22)
page_tag(s, 4, "")


# ================================================================ SLIDE 5 — SOURCE-BACKED / TRIALS
s = slide()
bg(s, WHITE)
rect(s, Inches(7.7), Inches(0), Inches(5.63), EMU_H, LIGHT)
kicker(s, Inches(0.9), Inches(0.7), "Evidence you can trust", TEAL)
text(s, Inches(0.9), Inches(1.05), Inches(6.4), Inches(1.6),
     [[("Every finding shows its work.", 32, INK, True, FONT)]],
     line_spacing=1.05)
text(s, Inches(0.9), Inches(2.35), Inches(6.3), Inches(3.2),
     [[("Source & identifier", 17, TEAL, True, FONT)],
      [("NCT IDs, PubMed citations, and direct links — never a black box.", 15, SLATE, False, FONT)],
      [("", 6, SLATE, False, FONT)],
      [("A reason it surfaced", 17, TEAL, True, FONT)],
      [("Biomarker alignment, recruitment status, geography, and freshness.", 15, SLATE, False, FONT)],
      [("", 6, SLATE, False, FONT)],
      [("Honest confidence labels", 17, TEAL, True, FONT)],
      [("High relevance · Worth reviewing · Low confidence · Insufficient data.", 15, SLATE, False, FONT)]],
     line_spacing=1.16)

# mock finding card on the right
cx, cy = Inches(8.15), Inches(1.1)
card(s, cx, cy, Inches(4.75), Inches(5.3), WHITE)
rect(s, cx, cy, Inches(4.75), Inches(0.16), TEAL_BR, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
chip(s, cx + Inches(0.3), cy + Inches(0.38), "High relevance", TEAL, WHITE)
chip(s, cx + Inches(2.05), cy + Inches(0.38), "Recruiting", TEAL_BR, NAVY)
text(s, cx + Inches(0.3), cy + Inches(0.92), Inches(4.2), Inches(1.0),
     [[("Targeted triple-combination for biomarker-positive disease",
        16, INK, True, FONT)]], line_spacing=1.1)
text(s, cx + Inches(0.3), cy + Inches(2.0), Inches(4.2), Inches(2.0),
     [[("Why it surfaced", 12, MUTED, True, FONT)],
      [("Matches subtype + biomarker; phase II; recruiting near your region.",
        13, SLATE, False, FONT)],
      [("", 5, SLATE, False, FONT)],
      [("Discuss with your team", 12, MUTED, True, FONT)],
      [("Is my variant eligible? What are the line-of-therapy requirements?",
        13, SLATE, False, FONT)]],
     line_spacing=1.16)
text(s, cx + Inches(0.3), cy + Inches(4.55), Inches(4.2), Inches(0.6),
     [[("Source:  ClinicalTrials.gov · NCT0000000   ↗", 12, TEAL, True, FONT)]])
page_tag(s, 5, "")


# ================================================================ SLIDE 6 — PRIVACY
s = slide()
bg(s, NAVY)
rect(s, Inches(0), Inches(0), Inches(0.32), EMU_H, CORAL)
kicker(s, Inches(0.9), Inches(0.7), "Private by design", CORAL)
text(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.4),
     [[("Identity stays on your device.", 34, WHITE, True, FONT)],
      [("Always your choice what leaves it.", 34, TEAL_BR, True, FONT)]],
     line_spacing=1.05)

# two mode cards
card(s, Inches(0.9), Inches(3.1), Inches(5.6), Inches(3.1), RGBColor(0x12,0x32,0x4D))
chip(s, Inches(1.2), Inches(3.4), "MODE 1", TEAL, WHITE)
text(s, Inches(1.2), Inches(3.95), Inches(5.0), Inches(0.6),
     [[("Local-only", 22, WHITE, True, FONT)]])
text(s, Inches(1.2), Inches(4.55), Inches(5.0), Inches(1.5),
     [[("Everything runs on-device. Public sources are queried with search terms only. "
        "No AI provider ever receives case context.", 15, RGBColor(0xC2,0xD2,0xDC), False, FONT)]],
     line_spacing=1.2)

card(s, Inches(6.8), Inches(3.1), Inches(5.6), Inches(3.1), RGBColor(0x12,0x32,0x4D))
chip(s, Inches(7.1), Inches(3.4), "MODE 2", TEAL_BR, NAVY)
text(s, Inches(7.1), Inches(3.95), Inches(5.0), Inches(0.6),
     [[("De-identified AI assist", 22, WHITE, True, FONT)]])
text(s, Inches(7.1), Inches(4.55), Inches(5.0), Inches(1.5),
     [[("Optional. Names, dates, and locations stay local. Only minimized, "
        "de-identified oncology context can be sent — after you acknowledge it.",
        15, RGBColor(0xC2,0xD2,0xDC), False, FONT)]],
     line_spacing=1.2)
rect(s, Inches(0.0), Inches(7.34), EMU_W, Inches(0.16), CORAL)


# ================================================================ SLIDE 7 — REPORTS
s = slide()
bg(s, CLOUD)
kicker(s, Inches(0.9), Inches(0.55), "Walk in prepared", TEAL)
text(s, Inches(0.9), Inches(0.86), Inches(11.5), Inches(0.9),
     [[("A report your oncologist can actually use.", 30, INK, True, FONT)]])

left = [
    "Patient profile snapshot",
    "New or changed items since last visit",
    "Possible trial matches",
    "Research & drug updates",
]
right = [
    "Why each item surfaced — and why it might not fit",
    "Missing information & data gaps",
    "Questions to discuss with your team",
    "Evidence appendix with every source",
]
card(s, Inches(0.9), Inches(2.0), Inches(5.6), Inches(4.4), WHITE)
card(s, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.4), WHITE)
text(s, Inches(1.2), Inches(2.25), Inches(5.0), Inches(0.5),
     [[("Daily Summary  &  Full Oncology Review", 16, TEAL, True, FONT)]])
y = Inches(2.95)
for item in left:
    rect(s, Inches(1.2), y + Inches(0.08), Inches(0.18), Inches(0.18), TEAL_BR, shape=MSO_SHAPE.OVAL)
    text(s, Inches(1.5), y - Inches(0.04), Inches(4.8), Inches(0.7),
         [[(item, 15.5, SLATE, False, FONT)]], line_spacing=1.1)
    y = y + Inches(0.82)
text(s, Inches(7.1), Inches(2.25), Inches(5.0), Inches(0.5),
     [[("Built to be reviewed by a clinician", 16, CORAL, True, FONT)]])
y = Inches(2.95)
for item in right:
    rect(s, Inches(7.1), y + Inches(0.08), Inches(0.18), Inches(0.18), CORAL, shape=MSO_SHAPE.OVAL)
    text(s, Inches(7.4), y - Inches(0.04), Inches(4.8), Inches(0.7),
         [[(item, 15.5, SLATE, False, FONT)]], line_spacing=1.1)
    y = y + Inches(0.82)
page_tag(s, 7, "")


# ================================================================ SLIDE 8 — HONEST SCOPE
s = slide()
bg(s, WHITE)
kicker(s, Inches(0.9), Inches(0.7), "Honest by default", TEAL)
text(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.9),
     [[("We tell you the truth about what this is.", 30, INK, True, FONT)]])

# IS card
card(s, Inches(0.9), Inches(2.2), Inches(5.6), Inches(4.0), LIGHT)
rect(s, Inches(0.9), Inches(2.2), Inches(0.16), Inches(4.0), TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(1.25), Inches(2.45), Inches(5.0), Inches(0.6),
     [[("Firstlight IS", 18, TEAL, True, FONT)]])
for i, t in enumerate(["An information monitor","A summarizer of public research",
                       "A preparation tool for visits","Transparent about its reasoning"]):
    text(s, Inches(1.25), Inches(3.1) + Inches(0.72)*i, Inches(5.0), Inches(0.6),
         [[("✓  ", 16, TEAL, True, FONT), (t, 15.5, SLATE, False, FONT)]])

# IS NOT card
card(s, Inches(6.8), Inches(2.2), Inches(5.6), Inches(4.0), RGBColor(0xFB,0xEF,0xEA))
rect(s, Inches(6.8), Inches(2.2), Inches(0.16), Inches(4.0), CORAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(7.15), Inches(2.45), Inches(5.0), Inches(0.6),
     [[("Firstlight is NOT", 18, CORAL, True, FONT)]])
for i, t in enumerate(["A diagnostic system","A treatment recommender",
                       "A trial-eligibility decision","A substitute for your oncologist"]):
    text(s, Inches(7.15), Inches(3.1) + Inches(0.72)*i, Inches(5.0), Inches(0.6),
         [[("—  ", 16, CORAL, True, FONT), (t, 15.5, SLATE, False, FONT)]])

text(s, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.6),
     [[("Every finding requires clinician review. Firstlight helps you ask better questions — not make medical decisions.",
        13.5, MUTED, False, FONT, True)]], align=PP_ALIGN.CENTER)
page_tag(s, 8, "")


# ================================================================ SLIDE 9 — WHO IT'S FOR
s = slide()
bg(s, NAVY)
rect(s, Inches(0), Inches(0), Inches(0.32), EMU_H, TEAL_BR)
kicker(s, Inches(0.9), Inches(0.7), "Who it's for", TEAL_BR)
text(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.9),
     [[("For the people doing the research at 2 a.m.", 32, WHITE, True, FONT)]])
who = [
    ("Patients", "Staying informed without drowning in journals."),
    ("Caregivers", "The spouse, the daughter, the friend who became the case manager."),
    ("Care teams", "Arriving to a structured, source-backed starting point."),
]
x = Inches(0.9)
for title, body in who:
    card(s, x, Inches(2.4), Inches(3.75), Inches(3.4), RGBColor(0x12,0x32,0x4D))
    rect(s, x + Inches(0.35), Inches(2.8), Inches(0.6), Inches(0.6), TEAL_BR, shape=MSO_SHAPE.OVAL)
    text(s, x + Inches(0.35), Inches(3.65), Inches(3.1), Inches(0.6),
         [[(title, 22, WHITE, True, FONT)]])
    text(s, x + Inches(0.35), Inches(4.3), Inches(3.1), Inches(1.4),
         [[(body, 15, RGBColor(0xC2,0xD2,0xDC), False, FONT)]], line_spacing=1.2)
    x = x + Inches(4.0)
rect(s, Inches(0.0), Inches(7.34), EMU_W, Inches(0.16), TEAL_BR)


# ================================================================ SLIDE 10 — VISION / ROADMAP
s = slide()
bg(s, CLOUD)
kicker(s, Inches(0.9), Inches(0.55), "Where it's going", TEAL)
text(s, Inches(0.9), Inches(0.86), Inches(11.5), Inches(0.9),
     [[("The watch keeps getting sharper.", 30, INK, True, FONT)]])
road = [
    ("Deeper trial feasibility", "Smarter geography and status filtering."),
    ("More safety connectors", "Additional drug label and safety sources."),
    ("Multi-profile households", "One app, more than one loved one."),
    ("Offline evidence caching", "Keep key sources available anywhere."),
    ("Stronger change diffing", "See exactly what moved since last time."),
    ("Auto-update channel", "Seamless improvements over time."),
]
cols = 3
cw = Inches(3.75); ch = Inches(1.85)
gx = Inches(0.9); gy = Inches(2.1)
for i, (t, b) in enumerate(road):
    r = i // cols; c = i % cols
    x = gx + (cw + Inches(0.13)) * c
    y = gy + (ch + Inches(0.2)) * r
    card(s, x, y, cw, ch, WHITE)
    rect(s, x + Inches(0.28), y + Inches(0.28), Inches(0.4), Inches(0.1), TEAL_BR, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.28), y + Inches(0.5), Inches(3.2), Inches(0.6),
         [[(t, 17, INK, True, FONT)]])
    text(s, x + Inches(0.28), y + Inches(1.05), Inches(3.25), Inches(0.7),
         [[(b, 13.5, SLATE, False, FONT)]], line_spacing=1.12)
page_tag(s, 10, "")


# ================================================================ SLIDE 11 — DEDICATION
s = slide()
bg(s, NAVY)
rect(s, Inches(0), Inches(0), EMU_W, Inches(0.16), CORAL)
rect(s, Inches(0), Inches(7.34), EMU_W, Inches(0.16), CORAL)
logo_mark(s, Inches(6.35), Inches(0.85), scale=1.2)
kicker(s, Inches(0), Inches(1.62), "Why it's called Firstlight", CORAL)
# center the kicker
s.shapes[-1].left = Inches(0); s.shapes[-1].width = EMU_W
s.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

text(s, Inches(1.6), Inches(2.1), Inches(10.13), Inches(1.7),
     [[("Firstlight is named for my mom.", 24, RGBColor(0xC9,0xD8,0xE0), True, SERIF)],
      [("", 6, WHITE, False, SERIF)],
      [("Judy Coffey", 46, AMBER, True, SERIF)]],
     align=PP_ALIGN.CENTER, line_spacing=1.05)

text(s, Inches(1.7), Inches(4.2), Inches(9.93), Inches(1.7),
     [[("We went searching through the research and found a promising trio of drugs from "
        "recent trials. We brought it to her doctor, and he agreed it was worth a try.",
        18, RGBColor(0xD8,0xE3,0xEA), False, SERIF, True)],
      [("", 5, WHITE, False, SERIF)],
      [("She passed before she could start it — in April 2025.", 18, CORAL, True, SERIF, True)]],
     align=PP_ALIGN.CENTER, line_spacing=1.22)

text(s, Inches(1.6), Inches(6.2), Inches(10.13), Inches(1.0),
     [[("So no other family finds what matters too late.", 18, WHITE, True, FONT)],
      [("In loving memory of Judy Coffey.  💛", 14, AMBER, True, FONT)]],
     align=PP_ALIGN.CENTER, line_spacing=1.25)


# ================================================================ SLIDE 12 — CTA
s = slide()
bg(s, NAVY)
rect(s, Inches(0), Inches(0), Inches(0.32), EMU_H, TEAL)
rect(s, Inches(0.32), Inches(0), Inches(0.10), EMU_H, TEAL_BR)
logo_mark(s, Inches(1.0), Inches(1.1), scale=1.1)
text(s, Inches(1.7), Inches(1.12), Inches(7), Inches(0.7),
     [[("Firstlight", 24, WHITE, True, FONT)]])
text(s, Inches(1.0), Inches(2.5), Inches(11), Inches(1.8),
     [[("Bring the latest research", 40, WHITE, True, FONT)],
      [("into the conversation that matters.", 40, TEAL_BR, True, FONT)]],
     line_spacing=1.05)
text(s, Inches(1.02), Inches(4.5), Inches(10.5), Inches(0.8),
     [[("Local-first.  Source-backed.  Private by design.  Free and open-source.",
        18, RGBColor(0xC9,0xD8,0xE0), False, FONT)]])

btn = rect(s, Inches(1.02), Inches(5.5), Inches(3.3), Inches(0.85), TEAL_BR, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tf = btn.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Download for macOS & Windows"; r.font.size = Pt(15)
r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = FONT

text(s, Inches(4.6), Inches(5.62), Inches(7), Inches(0.7),
     [[("github.com/seoplanotx/oncowatch", 16, WHITE, True, FONT)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(1.02), Inches(6.75), Inches(11), Inches(0.5),
     [[("Firstlight is an information tool, not medical advice. Every finding requires clinician review.",
        12, RGBColor(0x9B,0xB0,0xBE), False, FONT, True)]])
rect(s, Inches(0.0), Inches(7.34), EMU_W, Inches(0.16), TEAL_BR)


# ---------------------------------------------------------------- save
import os
out = os.path.join(os.path.dirname(__file__), "Firstlight-Promotional-Deck.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
